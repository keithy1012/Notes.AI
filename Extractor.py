from pptx import Presentation
import fitz  # PyMuPDF
import os
from PIL import Image
import pytesseract
import io

class Extractor:
    def __init__(self, path):
        self.path = path

    def extract_from_pptx(self):
        prs = Presentation(self.path)
        slides = []
        for i, slide in enumerate(prs.slides):
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text.strip())

            # Extract image and perform OCR
            if shape.shape_type == 13:
                image = shape.image
                image_bytes = image.blob
                try:
                    pil_image = Image.open(io.BytesIO(image_bytes))
                    ocr_text = pytesseract.image_to_string(pil_image)
                    if ocr_text.strip():
                        text.append(f"[OCR Text from Image]: {ocr_text.strip()}")
                except Exception as e:
                    print(f"⚠️ OCR failed on Slide {i+1}: {e}")

            slides.append({
                "section": f"Slide {i + 1}",
                "content": "\n".join(text)
            })
        return slides

    def extract_from_pdf(self):
        doc = fitz.open(self.path)
        pages = []
        for i, page in enumerate(doc):
            text = page.get_text()
            pages.append({
                "section": f"Page {i + 1}",
                "content": text.strip()
            })
        return pages

    def extract_from_txt(self):
        with open(self.path, "r", encoding="utf-8") as f:
            content = f.read()
        return [{"section": "Text Document", "content": content.strip()}]

    def extract_from_generic(self):
        try:
            with open(self.path, "r", encoding="utf-8") as f:
                content = f.read()
            return [{"section": "Generic File", "content": content.strip()}]
        except Exception as e:
            raise ValueError(f"Unsupported file type or unreadable content: {e}")

    def extract_content(self):
        ext = os.path.splitext(self.path)[1].lower()
        if ext == ".pptx":
            return self.extract_from_pptx()
        elif ext == ".pdf":
            return self.extract_from_pdf()
        elif ext == ".txt":
            return self.extract_from_txt()
        else:
            return self.extract_from_generic()